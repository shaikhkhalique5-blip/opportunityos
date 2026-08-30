import asyncio
from io import BytesIO

from fastapi import APIRouter, File, HTTPException, UploadFile
from pypdf import PdfReader
from pptx import Presentation

from app.core.db import SessionLocal
from app.models import DiscoveryRun
from app.schemas.discovery import DiscoveryRequest
from app.services.amplemarket import amplemarket_health, amplemarket_people_search_test
from app.services.discovery import push_to_amplemarket
from app.services.discovery_amplemarket import run_discovery

router = APIRouter(prefix="/discovery", tags=["discovery"])


def _set_run(run_id: int, **values):
    db = SessionLocal()
    try:
        row = db.get(DiscoveryRun, run_id)
        if row:
            for key, value in values.items():
                setattr(row, key, value)
            db.commit()
    finally:
        db.close()


def _progress_percent(stage: str) -> int:
    s = (stage or "").lower()
    if "complete" in s:
        return 100
    if "ranking" in s:
        return 72
    if "research" in s:
        return 46
    if "enrich" in s:
        return 38
    if "candidate" in s or "discover" in s:
        return 20
    if "product brain" in s:
        return 8
    if "queued" in s:
        return 2
    if "failed" in s:
        return 100
    return 5


async def _execute_run(run_id: int, req: DiscoveryRequest):
    _set_run(run_id, status="running", stage="Building Product Brain")

    def report(stage: str):
        _set_run(run_id, status="running", stage=stage)

    try:
        result = await run_discovery(req, progress=report)
        _set_run(run_id, status="completed", stage="Complete", response_json=result.model_dump(mode="json"), error=None)
    except Exception as exc:
        _set_run(run_id, status="failed", stage="Failed", error=f"{type(exc).__name__}: {str(exc)[:1500]}")


@router.post("/runs")
async def create_discovery_run(req: DiscoveryRequest):
    db = SessionLocal()
    try:
        row = DiscoveryRun(status="queued", stage="Queued", request_json=req.model_dump(mode="json"))
        db.add(row)
        db.commit()
        db.refresh(row)
        run_id = row.id
    finally:
        db.close()
    asyncio.create_task(_execute_run(run_id, req))
    return {"run_id": run_id, "status": "queued", "stage": "Queued", "progress": 2}


@router.get("/runs/{run_id}")
async def get_discovery_run(run_id: int):
    db = SessionLocal()
    try:
        row = db.get(DiscoveryRun, run_id)
        if not row:
            raise HTTPException(status_code=404, detail="Discovery run not found")
        return {
            "run_id": row.id,
            "status": row.status,
            "stage": row.stage,
            "progress": _progress_percent(row.stage),
            "result": row.response_json,
            "error": row.error,
        }
    finally:
        db.close()


# Compatibility endpoint. New UI uses /runs so browser requests never wait for deep research.
@router.post("/run")
async def discover_legacy(req: DiscoveryRequest):
    return await run_discovery(req)


@router.get("/amplemarket-health")
async def amplemarket_provider_health():
    return await amplemarket_health()


@router.get("/amplemarket-people-test")
async def amplemarket_people_provider_test():
    return await amplemarket_people_search_test()


@router.post("/deck")
async def ingest_deck(file: UploadFile = File(...)):
    try:
        raw = await file.read()
        name = (file.filename or "deck").lower()
        if len(raw) > 15 * 1024 * 1024:
            raise HTTPException(status_code=413, detail="Deck must be 15 MB or smaller")
        if name.endswith(".pdf"):
            reader = PdfReader(BytesIO(raw))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
        elif name.endswith(".pptx"):
            deck = Presentation(BytesIO(raw))
            parts = []
            for slide in deck.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            text = "\n".join(parts)
        elif name.endswith((".txt", ".md")):
            text = raw.decode("utf-8", errors="ignore")
        else:
            raise HTTPException(status_code=415, detail="Upload PDF, PPTX, TXT, or MD")
        text = text.strip()
        if not text:
            raise HTTPException(status_code=422, detail="No readable text found in deck")
        return {"filename": file.filename, "text": text[:60000], "characters": len(text)}
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not parse deck: {exc}") from exc


@router.post("/amplemarket")
async def amplemarket_handoff(payload: dict):
    try:
        return await push_to_amplemarket(payload)
    except Exception as exc:
        raise HTTPException(status_code=502, detail=str(exc)) from exc
